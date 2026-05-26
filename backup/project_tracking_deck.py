#!/usr/bin/env python3
"""
project_tracking_deck.py - 项目追踪表线上化方案演示文稿
基于 kami 设计系统 (parchment style)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ═══════════════════════════════════════════════════════════
# Design System Constants
# ═══════════════════════════════════════════════════════════

PARCHMENT   = RGBColor(0xf5, 0xf4, 0xed)
IVORY       = RGBColor(0xfa, 0xf9, 0xf5)
BRAND       = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK  = RGBColor(0x14, 0x14, 0x13)
DARK_WARM   = RGBColor(0x3d, 0x3d, 0x3a)
OLIVE       = RGBColor(0x50, 0x4e, 0x49)
STONE       = RGBColor(0x6b, 0x6a, 0x64)
BORDER      = RGBColor(0xe8, 0xe6, 0xdc)
WHITE       = RGBColor(0xff, 0xff, 0xff)

LANG = "zh"
CN_SERIF = "Source Han Serif SC"
CN_SANS  = "Microsoft YaHei" # Fallback sans-serif for body in PPT
SERIF = CN_SERIF
SANS  = CN_SANS

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def blank_slide(prs, bg_color=PARCHMENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide

def add_text(slide, text, left, top, width, height,
             font=SANS, size=18, bold=False, color=NEAR_BLACK,
             align=PP_ALIGN.LEFT, vanchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_rich_text(slide, left, top, width, height, items, align=PP_ALIGN.LEFT):
    """Add a textbox with multiple runs of different styles. items: [(text, font, size, bold, color)]"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    
    for i, (text, font, size, bold, color) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line

def add_card(slide, left, top, width, height, fill=IVORY, border=BORDER, border_weight=0.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_weight)
    card.shadow.inherit = False
    return card

# ═══════════════════════════════════════════════════════════
# Slide Generation
# ═══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    s = blank_slide(prs)
    add_text(s, "项目执行跟踪表线上化方案", Inches(1), Inches(2.2), Inches(11.33), Inches(1.5),
             font=SERIF, size=44, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.0), Inches(1), weight_pt=1.5)
    add_text(s, "多角色协同 · 审批闭环 · 动态管控", Inches(1), Inches(4.3), Inches(11.33), Inches(0.8),
             font=SANS, size=20, color=OLIVE, align=PP_ALIGN.CENTER)
    add_text(s, "Wood China 业务平台  ·  2026.05", Inches(1), Inches(6.5), Inches(11.33), Inches(0.4),
             font=SANS, size=13, color=STONE, align=PP_ALIGN.CENTER)

    # 2. TOC
    s = blank_slide(prs)
    add_text(s, "目录", Inches(1.2), Inches(0.8), Inches(10), Inches(0.8), font=SERIF, size=32, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(1.8), Inches(11), weight_pt=1)
    toc_items = ["业务痛点与核心目标", "角色分工与权限体系", "填报与审批流转闭环", "智能高亮与数据管控"]
    for i, item in enumerate(toc_items):
        y = Inches(2.4 + i * 0.9)
        add_text(s, f"0{i+1}", Inches(1.2), y, Inches(1), Inches(0.6), font=SERIF, size=28, color=BRAND)
        add_text(s, item, Inches(2.4), y, Inches(9), Inches(0.6), font=SERIF, size=22, color=NEAR_BLACK, vanchor=MSO_ANCHOR.MIDDLE)

    # 3. Chapter 1
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, "01", Inches(0.8), Inches(0.5), Inches(2), Inches(0.8), font=SERIF, size=26, color=WHITE)
    add_text(s, "业务痛点与核心目标", Inches(1), Inches(3), Inches(11.33), Inches(1.5), font=SERIF, size=56, color=WHITE, align=PP_ALIGN.CENTER)

    # 4. Pain Points
    s = blank_slide(prs)
    add_text(s, "核心目标", Inches(1.2), Inches(0.6), Inches(10), Inches(0.4), font=SANS, size=12, color=STONE)
    add_text(s, "从分散表格走向统一平台，消除数据孤岛", Inches(1.2), Inches(1.2), Inches(11), Inches(1), font=SERIF, size=30, color=NEAR_BLACK)
    
    # Cards
    pain_points = [
        ("数据不一致", "CRB 合同额与线下 Excel 手工填报存在差异，多版本导致信任危机。"),
        ("权限无管控", "全员可编辑导致误操作，缺乏细粒度的字段级权限保护。"),
        ("追溯困难", "历史数据覆盖后无法找回，缺乏审批快照与审计日志。"),
        ("协同低效", "邮件反复收发 Excel，版本汇总耗时耗力。")
    ]
    for i, (title, desc) in enumerate(pain_points):
        col = i % 2
        row = i // 2
        x = Inches(1.2 + col * 5.6)
        y = Inches(2.8 + row * 2.2)
        add_card(s, x, y, Inches(5.2), Inches(1.8))
        add_text(s, title, x + Inches(0.3), y + Inches(0.2), Inches(4.5), Inches(0.5), font=SERIF, size=20, bold=True, color=BRAND)
        add_text(s, desc, x + Inches(0.3), y + Inches(0.8), Inches(4.5), Inches(0.8), font=SANS, size=14, color=DARK_WARM)

    # 5. Chapter 2
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, "02", Inches(0.8), Inches(0.5), Inches(2), Inches(0.8), font=SERIF, size=26, color=WHITE)
    add_text(s, "角色分工与权限体系", Inches(1), Inches(3), Inches(11.33), Inches(1.5), font=SERIF, size=56, color=WHITE, align=PP_ALIGN.CENTER)

    # 6. Roles
    s = blank_slide(prs)
    add_text(s, "角色矩阵", Inches(1.2), Inches(0.6), Inches(10), Inches(0.4), font=SANS, size=12, color=STONE)
    add_text(s, "六大角色协同，实现业务填报与财务核算权责分离", Inches(1.2), Inches(1.2), Inches(11), Inches(1), font=SERIF, size=30, color=NEAR_BLACK)
    
    roles = [
        ("🛡️ 系统管理员", "系统运维、特批解锁、报表配置"),
        ("💰 财务审核", "每月 1-3 日专属，全量核对开票/回款"),
        ("👥 板块管理员", "板块数据汇总、异常处理、CRB 映射"),
        ("👷 项目经理", "填报产值、进度、WIP，支持导入导出"),
        ("📊 板块总监", "查看板块数据，一级审批确认"),
        ("🏛️ 项目群群主", "跨板块查看，二级审批确认"),
    ]
    for i, (role, desc) in enumerate(roles):
        y = Inches(2.8 + i * 0.7)
        add_text(s, role, Inches(1.2), y, Inches(2.5), Inches(0.6), font=SERIF, size=18, bold=True, color=BRAND)
        add_text(s, desc, Inches(4.0), y, Inches(8.5), Inches(0.6), font=SANS, size=16, color=DARK_WARM)

    # 7. Chapter 3
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, "03", Inches(0.8), Inches(0.5), Inches(2), Inches(0.8), font=SERIF, size=26, color=WHITE)
    add_text(s, "填报与审批流转闭环", Inches(1), Inches(3), Inches(11.33), Inches(1.5), font=SERIF, size=56, color=WHITE, align=PP_ALIGN.CENTER)

    # 8. Timeline
    s = blank_slide(prs)
    add_text(s, "月度作业流", Inches(1.2), Inches(0.6), Inches(10), Inches(0.4), font=SANS, size=12, color=STONE)
    add_text(s, "以 1/19/25/9 为关键节点的周期性管控", Inches(1.2), Inches(1.2), Inches(11), Inches(1), font=SERIF, size=30, color=NEAR_BLACK)
    
    timeline = [
        ("1日 - 3日", "财务专属期\n全量核对开票、回款、预收款数据", BRAND),
        ("19日", "填报提醒\n自动发送消息督促各角色完成填写", STONE),
        ("25日", "月度锁定\n填报窗口关闭，除管理员外全员冻结", BRAND),
        ("次月 9日", "常规解禁\n解锁当月修正与未来月份预测", OLIVE)
    ]
    for i, (date, desc, col) in enumerate(timeline):
        x = Inches(1.2 + i * 3.0)
        add_card(s, x, Inches(2.8), Inches(2.6), Inches(3.5))
        add_text(s, date, x + Inches(0.2), Inches(3.0), Inches(2.2), Inches(0.5), font=SERIF, size=22, bold=True, color=col)
        add_text(s, desc, x + Inches(0.2), Inches(3.7), Inches(2.2), Inches(2.0), font=SANS, size=14, color=DARK_WARM)

    # 9. Approval
    s = blank_slide(prs)
    add_text(s, "审批快照流", Inches(1.2), Inches(0.6), Inches(10), Inches(0.4), font=SANS, size=12, color=STONE)
    add_text(s, "四级流转机制，确保数据全程可追溯", Inches(1.2), Inches(1.2), Inches(11), Inches(1), font=SERIF, size=30, color=NEAR_BLACK)
    
    steps = [
        ("Draft", "草稿版", "板块/中心完成填报，系统锁定生成快照"),
        ("Approve1", "总监版", "板块总监审核通过，生成二级快照"),
        ("Approve2", "群主版", "项目群群主跨板块核对确认"),
        ("J 版", "归档版", "系统管理员最终确认，正式归档封存"),
    ]
    for i, (code, name, desc) in enumerate(steps):
        y = Inches(2.8 + i * 1.0)
        add_text(s, code, Inches(1.2), y, Inches(1.5), Inches(0.6), font=SERIF, size=24, bold=True, color=BRAND)
        add_text(s, name, Inches(2.8), y, Inches(1.5), Inches(0.6), font=SERIF, size=20, color=NEAR_BLACK, vanchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, Inches(4.5), y, Inches(8), Inches(0.6), font=SANS, size=16, color=DARK_WARM, vanchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            add_line(s, Inches(1.8), y + Inches(0.7), Inches(0.4), weight_pt=2)

    # 10. Chapter 4
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, "04", Inches(0.8), Inches(0.5), Inches(2), Inches(0.8), font=SERIF, size=26, color=WHITE)
    add_text(s, "智能高亮与数据管控", Inches(1), Inches(3), Inches(11.33), Inches(1.5), font=SERIF, size=56, color=WHITE, align=PP_ALIGN.CENTER)

    # 11. Highlights & Logic
    s = blank_slide(prs)
    add_text(s, "智能管控", Inches(1.2), Inches(0.6), Inches(10), Inches(0.4), font=SANS, size=12, color=STONE)
    add_text(s, "Diff 高亮与自动初始化，让变更一目了然", Inches(1.2), Inches(1.2), Inches(11), Inches(1), font=SERIF, size=30, color=NEAR_BLACK)

    # Left side: Init
    add_card(s, Inches(1.2), Inches(2.8), Inches(5.2), Inches(3.5))
    add_text(s, "🔄 数据初始化规则", Inches(1.4), Inches(3.0), Inches(4.8), Inches(0.5), font=SERIF, size=22, bold=True, color=BRAND)
    init_items = ["当月完成/开票/回款自动归零", "历史月份数据只读保留", "未来预测数据继承并允许修改"]
    for i, item in enumerate(init_items):
        add_text(s, f"• {item}", Inches(1.6), Inches(3.8 + i*0.7), Inches(4.5), Inches(0.5), font=SANS, size=16, color=DARK_WARM)

    # Right side: Highlights
    add_card(s, Inches(7.0), Inches(2.8), Inches(5.2), Inches(3.5))
    add_text(s, "✨ 变更高亮机制", Inches(7.2), Inches(3.0), Inches(4.8), Inches(0.5), font=SERIF, size=22, bold=True, color=BRAND)
    hl_items = ["新增项目整行浅色高亮", "未来预测修改触发单元格高亮", "联动计算字段受牵连同步高亮"]
    for i, item in enumerate(hl_items):
        add_text(s, f"• {item}", Inches(7.4), Inches(3.8 + i*0.7), Inches(4.5), Inches(0.5), font=SANS, size=16, color=DARK_WARM)

    # 12. Ending
    s = blank_slide(prs)
    add_text(s, "感谢观看", Inches(1), Inches(3), Inches(11.33), Inches(1.2), font=SERIF, size=40, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.5), Inches(1), weight_pt=1.5)
    add_text(s, "期待与您共同推进业务数字化", Inches(1), Inches(4.8), Inches(11.33), Inches(0.6), font=SANS, size=16, color=OLIVE, align=PP_ALIGN.CENTER)

    output_path = r"C:\Work\1_Projects\202605_项目追踪表线上化\项目追踪表线上化方案.pptx"
    prs.save(output_path)
    print(f"OK: Saved {output_path}")

if __name__ == '__main__':
    main()
